import json
from datetime import date
from pathlib import Path
from typing import Dict, List, Optional

from rich.console import Console

console = Console()


def export_run(
    generated: Dict,
    synthesis: Dict,
    manifest: Dict,
    output_dir: Path,
) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    image_prompts: List[str] = []

    for fmt_name, content_data in generated.items():
        if fmt_name == "article":
            _export_article(content_data, output_dir)
        elif fmt_name == "carousel":
            prompts = _export_carousel(content_data, output_dir)
            image_prompts.extend(prompts)
        elif fmt_name == "infographic":
            prompt = _export_infographic(content_data, output_dir)
            if prompt:
                image_prompts.append(prompt)
        elif fmt_name == "short_post":
            _export_short_post(content_data, output_dir)

    if image_prompts:
        header = "# Image Generation Prompts\n\nUse these with ChatGPT, DALL-E, or Midjourney.\n\n"
        (output_dir / "image_prompts.md").write_text(
            header + "\n\n---\n\n".join(image_prompts)
        )
        console.print(f"  Image prompts → [dim]{output_dir / 'image_prompts.md'}[/dim]")

    summary = {
        "date": date.today().isoformat(),
        "conference": manifest.get("conference_name", ""),
        "formats_generated": list(generated.keys()),
        "themes": [t["title"] for t in synthesis.get("themes", [])],
        "hot_takes_count": len(synthesis.get("hot_takes", [])),
    }
    (output_dir / "run_summary.json").write_text(json.dumps(summary, indent=2))
    console.print(f"  [green]All output → {output_dir}[/green]")


def _export_article(data: Dict, output_dir: Path) -> None:
    path = output_dir / "article.md"
    path.write_text(data["content"])
    console.print(f"  Article → [dim]{path}[/dim]")


def _export_carousel(data: Dict, output_dir: Path) -> List[str]:
    (output_dir / "carousel.md").write_text(data["content"])
    console.print(f"  Carousel (text) → [dim]{output_dir / 'carousel.md'}[/dim]")

    slides_data: Optional[Dict] = data.get("slides_data")
    image_prompts: List[str] = []

    if slides_data and slides_data.get("slides"):
        try:
            pptx_path = output_dir / "carousel.pptx"
            _create_carousel_pptx(slides_data["slides"], pptx_path)
            console.print(f"  Carousel (PPTX) → [dim]{pptx_path}[/dim]")
        except Exception as e:
            console.print(f"  [yellow]PPTX generation failed: {e}[/yellow]")

        for slide in slides_data.get("slides", []):
            prompt = slide.get("image_prompt", "")
            if prompt:
                headline = slide.get("headline", f"Slide {slide.get('slide_number', '')}")
                image_prompts.append(f"## Carousel — {headline}\n\n{prompt}")

    return image_prompts


def _export_infographic(data: Dict, output_dir: Path) -> Optional[str]:
    path = output_dir / "infographic.md"
    path.write_text(data["content"])
    console.print(f"  Infographic → [dim]{path}[/dim]")

    content = data["content"]
    if "IMAGE GENERATION PROMPT" in content.upper() or "DALL-E" in content.upper():
        lines = content.splitlines()
        collecting = False
        prompt_lines = ["## Infographic — Image Generation Prompt"]
        for line in lines:
            upper = line.upper().strip()
            if any(k in upper for k in ["IMAGE GENERATION PROMPT", "DALL-E PROMPT", "DALL-E/"]):
                collecting = True
                continue
            # Stop when the next ## section starts (after we've begun collecting)
            if collecting and upper.startswith("##"):
                break
            if collecting:
                prompt_lines.append(line)
        if len(prompt_lines) > 1:
            return "\n".join(prompt_lines)

    return f"## Infographic\n\n{content[:500]}"


def _export_short_post(data: Dict, output_dir: Path) -> None:
    path = output_dir / "short_post.md"
    path.write_text(data["content"])
    console.print(f"  Short post → [dim]{path}[/dim]")


# ── PPTX Carousel ────────────────────────────────────────────────────────────

def _create_carousel_pptx(slides: List[Dict], output_path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    DARK_BG = RGBColor(0x0D, 0x11, 0x17)
    LINKEDIN_BLUE = RGBColor(0x00, 0x77, 0xB5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x8B, 0x94, 0x9E)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10)
    blank = prs.slide_layouts[6]

    def bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def accent_bar(slide):
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = LINKEDIN_BLUE
        bar.line.fill.background()

    def textbox(slide, text, left, top, width, height, size,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def slide_num_badge(slide, num):
        badge = slide.shapes.add_shape(
            1, Inches(0.45), Inches(0.45), Inches(0.65), Inches(0.65)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = LINKEDIN_BLUE
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(num)
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for slide_data in slides:
        slide = prs.slides.add_slide(blank)
        bg(slide)
        stype = slide_data.get("type", "content")

        if stype == "title":
            accent_bar(slide)
            textbox(slide, slide_data.get("headline", ""),
                    0.7, 2.2, 8.6, 4.5, 40, bold=True, align=PP_ALIGN.CENTER)
            if slide_data.get("subtitle"):
                textbox(slide, slide_data["subtitle"],
                        0.7, 7.2, 8.6, 1.2, 18, color=GRAY, align=PP_ALIGN.CENTER)

        elif stype == "content":
            num = slide_data.get("slide_number")
            if num:
                slide_num_badge(slide, num)
            textbox(slide, slide_data.get("headline", ""),
                    0.5, 1.5, 9.0, 1.8, 30, bold=True)
            bullets = slide_data.get("bullets", [])
            if bullets:
                textbox(slide, "\n".join(f"• {b}" for b in bullets),
                        0.5, 3.6, 9.0, 5.5, 22, color=GRAY)

        elif stype == "stat":
            textbox(slide, slide_data.get("stat", ""),
                    0.5, 2.0, 9.0, 3.0, 96, bold=True,
                    color=LINKEDIN_BLUE, align=PP_ALIGN.CENTER)
            textbox(slide, slide_data.get("context", ""),
                    0.8, 5.5, 8.4, 2.0, 24, align=PP_ALIGN.CENTER)
            if slide_data.get("source"):
                textbox(slide, f"— {slide_data['source']}",
                        0.8, 7.8, 8.4, 0.8, 14, color=GRAY, align=PP_ALIGN.CENTER)

        elif stype == "quote":
            textbox(slide, "“", 0.5, 0.6, 2.5, 2.0, 96,
                    bold=True, color=LINKEDIN_BLUE)
            textbox(slide, slide_data.get("quote", ""),
                    0.9, 2.5, 8.2, 4.5, 26)
            if slide_data.get("attribution"):
                textbox(slide, f"— {slide_data['attribution']}",
                        0.9, 7.5, 8.2, 1.0, 16, color=GRAY, align=PP_ALIGN.RIGHT)

        elif stype == "cta":
            accent_bar(slide)
            textbox(slide, slide_data.get("headline", ""),
                    0.7, 2.8, 8.6, 2.8, 36, bold=True, align=PP_ALIGN.CENTER)
            if slide_data.get("sub"):
                textbox(slide, slide_data["sub"],
                        0.7, 6.2, 8.6, 2.5, 20, color=GRAY, align=PP_ALIGN.CENTER)

    prs.save(str(output_path))
